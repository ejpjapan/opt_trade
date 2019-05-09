

from pathlib import Path
import pandas as pd
from pptx import Presentation
from spx_data_update import UpdateSP500Data
from option_utilities import read_feather, write_feather
from urllib.request import urlretrieve
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
import pandas_datareader.data as web
import os


def title_only_slide(asset_path, slide_dict, presentation, text_msg=None):
    slide = presentation.slides.add_slide(presentation.slide_layouts[slide_dict['TITLE_ONLY']])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    placeholder = slide.placeholders[13]  # idx key, not position
    placeholder_picture = placeholder.insert_picture(str(asset_path))
    slide.shapes.title.text = text_msg
    return slide


def main():
    ppt_path = Path.home() / 'Dropbox' / 'option_overlay'
    fig_path = Path.home() / 'Dropbox' / 'outputDev' / 'fig'
    template_name = 'option_overlay.pptx'
    output_name = 'test.pptx'

    # Assets
    heat_map_path = fig_path / 'heat_map.png'
    cum_perf_path = fig_path / 'cum_perf.png'
    cum_total_perf_path = fig_path / 'cum_total_perf.png'
    cum_total_perf2_path = fig_path / 'cum_total_perf2.png'
    hfri_overlay_path = fig_path / 'hfri_overlay.png'
    delta_path = fig_path / 'Delta.png'
    gamma_path = fig_path / 'Gamma.png'
    vega_path = fig_path / 'Vega.png'
    rho_path = fig_path / 'Rho.png'
    theta_path = fig_path / 'Theta.png'

    # Layout index
    layout_dict = {'TITLE': 0, 'SUB_TITLE': 1, 'QUOTE': 2, 'TITLE_COLUMN1': 3, 'TITLE_COLUMN2': 4, 'TITLE_COLUMN3': 5,
                   'TITLE_ONLY': 6, 'CAPTION': 7, 'BLANK': 8}

    prs = Presentation(ppt_path / template_name)

    # Title slide
    for shape in prs.slides[0].placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    prs.slides[0].shapes[0].text = 'Income Alternatives'

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_dict['TITLE_COLUMN1']])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    # placeholder = slide.placeholders[1]  # idx key, not position
    slide.shapes.title.text = 'Background'

    paragraph_strs = [
        'Egg, bacon, sausage and spam.',
        'Spam, bacon, sausage and spam.',
        'Spam, egg, spam, spam, bacon and spam.'
    ]
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()  # remove any existing paragraphs, leaving one empty one
    
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    p.alignment = PP_ALIGN.LEFT

    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # 2
    title_only_slide(heat_map_path, layout_dict, prs, text_msg='Monthly Excess Returns (%)')
    # 3
    title_only_slide(cum_perf_path, layout_dict, prs, text_msg='Cumulative Excess Return')
    # 4
    title_only_slide(cum_total_perf_path, layout_dict, prs, text_msg='Cumulative Total Return')
    # 5
    title_only_slide(cum_total_perf2_path, layout_dict, prs, text_msg='Cumulative Total Return')
    # 6
    title_only_slide(hfri_overlay_path, layout_dict, prs, text_msg='Overlay vs. HFRI')

    # 6~10
    greek_dict = {delta_path: 'Delta',
                  gamma_path: 'Gamma',
                  vega_path: 'Vega',
                  rho_path: 'Rho',
                  theta_path: 'Theta'}
    for key, value in greek_dict.items():
        title_only_slide(key, layout_dict, prs, text_msg='Portfolio ' + value)

    # Save and open presentation
    prs.save(ppt_path / output_name)
    os.system("open " + str(ppt_path / output_name))

    #
    # greek_dict = {delta_path: 'Delta',
    # gamma_path: 'Gamma',
    # vega_path: 'Vega',
    # rho_path: 'Rho',
    # theta_path: 'Theta.png'}
    # for key, value in greek_dict.items():
    #     slide = prs.slides.add_slide(prs.slide_layouts[layout_dict['TITLE_ONLY']])


if __name__ == '__main__':
    main()



# for i in range(0, 8, 1):
#     blank_slide_layout = prs.slide_layouts[i]
#     slide = prs.slides.add_slide(blank_slide_layout)
#
# top = Inches(1.54)
# left = Inches(0.28)
# height = Inches(3.82)
# pic = slide.shapes.add_picture(str(heat_map_path), left, top, height=height)


# for shape in slide.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.name))

def aqr_alt_funds(update_funds=True):

    db_directory = UpdateSP500Data.DATA_BASE_PATH / 'xl'
    url_string = 'https://funds.aqr.com/-/media/files/fund-documents/pricefiles/'

    fund_dict = {'alternative_risk_premia': 'leapmf.xls',
                 'diversified_arbitrage': 'daf.xls',
                 'equity_market_neutral': 'emnmf.xls',
                 'equity_long_short': 'elsmf.xls',
                 'global_macro': 'gmmf.xls',
                 'managed_futures': 'mfmf.xls',
                 'multi_alternative': 'msaf.xls',
                 'style_premia_alternative': 'spaf.xls'}

    url_dict = {value: url_string + value for (key, value) in fund_dict.items()}

    if update_funds:
        _ = [urlretrieve(value, db_directory / key) for (key, value) in url_dict.items()]

    rows_to_skip = list(range(0, 15))
    rows_to_skip.append(16)

    aqr_funds_index = []
    for key, value in fund_dict.items():
        df = pd.read_excel(db_directory / value, usecols=[1, 4],
                           skiprows=rows_to_skip, index_col=0, squeeze=True,
                           keep_default_na=False)
        df = df.rename(key)
        aqr_funds_index.append(df)
    return pd.concat(aqr_funds_index, axis=1)


def get_fund_assets(update_funds=True):

    db_directory = UpdateSP500Data.DATA_BASE_PATH / 'feather'
    feather_name = 'all_funds.feather'
    if update_funds:
        fund_dict = {'^SP500TR': 'S&P 500',
            'VDIGX': 'VG Dividend Growth',
            'VEIRX': 'VG Equity-Income',
            'VWEAX': 'VG High-Yield Corporate',
            'VWALX': 'VG High-Yield Muni',
            'VBTLX': 'VG Total Bond Market',
            'BXMIX': 'Blackstone Alternatives',
            'QLEIX': 'AQR Equity Long/Short',
            'QGMIX': 'AQR Global Macro',
            'QMHIX': 'AQR Managed Futures',
            'ADAIX': 'AQR Diversified Arbitrage',
            'QSPIX': 'AQR Style Premia',
            'AVGRX': 'Dreyfus Dynamic Total Return', #$1.141bn
            'FAAAX': 'K2 Franklin Alternative',# fund $1.17bn
            'GJRTX': 'GSAM Absolute return', # tracker $2.36bn
            'MASNX': 'Litman Gregory Masters Alt',# Strats Fund $2.05bn
            'PSMIX': 'Principal Global Multi-Strategy',# Fund $2.76bn
            'QOPIX': 'Oppenheimer Fundamental Alternatives',# Fd $1.20
            'GAFYX': 'Natixis ASG Global Alternatives'} #  Fd $1.39bn

        all_funds = [web.get_data_yahoo(key, 'JAN-16-80') for key, _ in fund_dict.items()]
        all_funds = [fund['Adj Close'] for fund in all_funds]
        all_funds = [fund.rename(fund_name) for fund, fund_name in zip(all_funds, fund_dict.values())]
        all_funds = pd.concat(all_funds, axis=1)
        # Replace dodgy observation
        all_funds['Principal Global Multi-Strategy']['2017-08-24'] = all_funds['Principal Global Multi-Strategy'][
            '2017-08-23']
        write_feather(all_funds, db_directory / feather_name)

    all_funds = read_feather(db_directory / feather_name)
    return all_funds


def get_hfr(feather_name, csv_file_path, update_funds=True):
    db_directory = UpdateSP500Data.DATA_BASE_PATH / 'xl'
    # feather_name = 'hfrx.feather'
    if update_funds:
        rows_to_skip = list(range(0, 2))
        headers = ['Date', 'Index Name', 'Index Code', 'Return', 'Index Value']

        df = pd.read_csv(db_directory / csv_file_path, skiprows=rows_to_skip,
                         squeeze=True, names=headers, engine='python')
        index_codes = df['Index Code'].unique()
        all_hfr_list = []
        for index_code in index_codes[:-1]: # remove HFR company info
            idx = df['Index Code'] == index_code
            hfr = df[idx].copy()
            hfr['Date'] = hfr['Date'].apply(pd.to_datetime)
            hfr = hfr.set_index(['Date'])
            hfr = hfr.reindex(hfr.index.sort_values())
            hfr_index = hfr['Index Value'].rename(hfr['Index Name'].unique()[0])
            all_hfr_list.append(hfr_index)
        hfr_df = pd.concat(all_hfr_list, axis=1)
        write_feather(hfr_df, UpdateSP500Data.DATA_BASE_PATH / 'feather' / feather_name)

    hfr_df = read_feather(UpdateSP500Data.DATA_BASE_PATH / 'feather' / feather_name)
    return hfr_df


def daily_hfrx(**kwargs):
    feather_name = 'hfrx.feather'
    csv_file_path = 'hfrx_daily_index_data.csv'
    return get_hfr(feather_name, csv_file_path, **kwargs)


def monthly_hfri(**kwargs):
    feather_name = 'hfri.feather'
    csv_file_path = 'hfri_index_data.csv'
    return get_hfr(feather_name, csv_file_path, **kwargs)


def monthly_hfrr(**kwargs):
    feather_name = 'hfrr.feather'
    csv_file_path = 'hfrr_index_data.csv'
    return get_hfr(feather_name, csv_file_path, **kwargs)

# def daily_hfrx(update_funds=True):
#     db_directory = UpdateSP500Data.DATA_BASE_PATH / 'xl'
#     feather_name = 'hfrx.feather'
#     if update_funds:
#         rows_to_skip = list(range(0, 2))
#         headers = ['Date', 'Index Name', 'Index Code', 'Return', 'Index Value']
#
#         df = pd.read_csv(db_directory / 'hfrx_daily_index_data.csv', skiprows=rows_to_skip,
#                          squeeze=True, names=headers, engine='python')
#         index_codes = df['Index Code'].unique()
#         all_hfrx_list = []
#         for index_code in index_codes[:-1]: # remove HFR company info
#             idx = df['Index Code'] == index_code
#             hfr = df[idx]
#             hfr.loc[:, 'Date'] = pd.to_datetime(hfr.loc[:, 'Date'])
#             hfr = hfr.set_index(['Date'])
#             hfr = hfr.reindex(hfr.index.sort_values())
#             hfr_index = hfr['Index Value'].rename(hfr['Index Name'].unique()[0])
#             all_hfrx_list.append(hfr_index)
#         hfrx = pd.concat(all_hfrx_list, axis=1)
#         write_feather(hfrx, UpdateSP500Data.DATA_BASE_PATH / 'feather' / feather_name)
#
#     hfrx = read_feather(UpdateSP500Data.DATA_BASE_PATH / 'feather' / feather_name)
#     return hfrx


